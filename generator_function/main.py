from sqlalchemy import create_engine
from sqlalchemy.orm import sessionmaker
# импортируем классы Book и Base из файла database_setup.py
from database_setup import Users
from database_setup import Base
engine = create_engine('mysql+pymysql://remote:Password123!!!@40.115.108.198/garazh_analytics')# Свяжим engine с метаданными класса Base,
# чтобы декларативы могли получить доступ через экземпляр DBSession
Base.metadata.bind = engine

DBSession = sessionmaker(bind=engine)
# Экземпляр DBSession() отвечает за все обращения к базе данных
# и представляет «промежуточную зону» для всех объектов,
# загруженных в объект сессии базы данных.
session = DBSession()
bookOne = Users(username="hui", role="1")
session.add(bookOne)
session.commit()

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import ChartData
from pptx.util import Inches
import numpy as np
import datetime

title='Automated Presentation Creating Process\n\
How to Create PowerPoint Presentations with Python '
prs = Presentation()

slide = prs.slides.add_slide(prs.slide_layouts[6])
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(9/1.5),Inches(16),Inches(9/8.5)
)
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(255,0,0)
shape.text= title
line=shape.line
line.color.rgb=RGBColor(255,0,0)

slide = prs.slides.add_slide(prs.slide_layouts[6])

shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.3))
shape.shadow.inherit = False
fill=shape.fill
fill.solid()
fill.fore_color.rgb=RGBColor(255,0,0)
shape.text= "How to Add a Chart"
line=shape.line
line.color.rgb=RGBColor(255,0,0)


N = 100

random_x = np.random.randn(N) + 10
random_y = np.random.randn(N)+5
random_z = np.random.randn(N) +20

dte=datetime.datetime.today()
dt_lst=[dte-datetime.timedelta(days=i) for i in range(N)]

chart_data = ChartData()
chart_data.categories = dt_lst
chart_data.add_series('Data 1',    random_x)
chart_data.add_series('Data 2',    random_y)
chart_data.add_series('Data 3',    random_z)

x, y, cx, cy = Inches(1), Inches(2), Inches(14), Inches(6)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
).chart
chart.has_legend = True
chart.legend.include_in_layout = False
chart.series[2].smooth = True

prs.save('basic_presentation_2.pptx')