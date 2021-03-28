#! /usr/bin/env python
# -*- coding: utf-8 -*-

from typeform import Typeform
from flask import Flask
import json
import plotly.express as px
import plotly.graph_objects as go
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import ChartData
from pptx.util import Inches
import numpy as np
import datetime
import pandas as pd

def get_typeform_data():
    # обрабатываем ответы
    responses = Typeform('FfLZwidfYeoAbHjamv3QCUpDimgA9oXFg6KUWhBS6yyz').responses
    result = responses.list('or3ePC3s')

    result = str(result).replace("'", '"')

    json_object = json.dumps(result)

    f = open("main.json", "w")
    f.write(result)
    f.close()

    f = open('main.json')
    data = json.load(f)

    dict_array = []
    for j in range(len(data['items'])):
      answers = []
      questions_id = []
      for i in range(len(data['items'][j]['answers'])):
        questions_id.append(data['items'][j]['answers'][i]['field']['id'])
        if data['items'][j]['answers'][i]['type'] == 'choice':
          answers.append(data['items'][j]['answers'][i]['choice']['label'])
        if data['items'][j]['answers'][i]['type'] == 'text' or data['items'][j]['answers'][i]['type'] == 'long_text':
          answers.append(data['items'][j]['answers'][i]['text'])
      new_dict = dict(zip(questions_id, answers))
      dict_array.append(new_dict)

    # обрабатываем вопросы

    forms = Typeform('FfLZwidfYeoAbHjamv3QCUpDimgA9oXFg6KUWhBS6yyz').forms
    result = forms.get('or3ePC3s')

    result = str(result).replace("'", '"')
    result = str(result).replace('True', '"True"')
    result = str(result).replace('False', '"False"')

    json_object = json.dumps(result)

    f = open("main.json", "w")
    f.write(result)
    f.close()

    f = open('main.json')
    data = json.load(f)

    id = []
    title = []
    for i in range(len(data['fields'])):
      id.append(data['fields'][i]['id'])
      title.append(data['fields'][i]['title'])

    question_titles = dict(zip(id, title))

    temp_dict = dict_array[0]
    a = pd.DataFrame.from_dict(temp_dict, orient='index')
    b = pd.DataFrame.from_dict(question_titles, orient='index')
    final_table = a.merge(b, left_index=True, right_index=True)
    for i in range(1, len(dict_array)):
        temp_dict = dict_array[i]
        a = pd.DataFrame.from_dict(temp_dict, orient='index')
        temp = a.merge(b, left_index=True, right_index=True)
        final_table = pd.concat([final_table, temp])

    final_table.columns = ['answer', 'question']
    return final_table

def generate_pptx_file():

    df = get_typeform_data()
    # пример круговой диаграммы

    tmp = df.loc[pd.Index(['Hx0Y0ufJuFB1'])].groupby(['answer']).count().reset_index()

    fig = go.Figure(data=[go.Pie(labels=tmp['answer'].tolist(), values=tmp['question'].tolist())])
    fig.write_image("img.png")

    # пример гистограммы

    tmp = df.loc[pd.Index(['l7AsgL0PJ32Q'])].groupby(['answer']).count().reset_index()

    dt = px.data.tips()
    fig = px.histogram(tmp, x="answer")
    fig.show()
    fig.write_image("hist.png")

    # Генерируем презентацию

    title='Automated Presentation Creating Process\n\
    Garazh Moscow Travel Hack '
    prs = Presentation()

    # front page
    #-----------------------------------------------------------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)



    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(9/1.5),Inches(16),Inches(9/8.5)
    )
    shape.shadow.inherit = False
    fill=shape.fill
    fill.solid()
    fill.fore_color.rgb=RGBColor(255,0,0)
    shape.text= title
    line=shape.line
    line.color.rgb=RGBColor(255,0,0)
    #-----------------------------------------------------------------------------------------------------------------------

    #Page 1
    #-----------------------------------------------------------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.3))
    shape.shadow.inherit = False
    fill=shape.fill
    fill.solid()
    fill.fore_color.rgb=RGBColor(255,0,0)
    shape.text= "People analytics"
    line=shape.line
    line.color.rgb=RGBColor(255,0,0)

    N = 100

    random_x = np.random.randn(N) + 10
    random_y = np.random.randn(N)+5
    random_z = np.random.randn(N) +20

    dte=datetime.datetime.today()
    dt_lst=[dte-datetime.timedelta(days=i) for i in range(N)]

    chart_data = ChartData()
    chart_data.categories = dt_lst
    chart_data.add_series('Data 1',    random_x)
    chart_data.add_series('Data 2',    random_y)
    chart_data.add_series('Data 3',    random_z)

    x, y, cx, cy = Inches(1), Inches(2), Inches(14), Inches(6)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    ).chart
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.series[2].smooth = True


    #Page 2
    #-----------------------------------------------------------------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.3))
    shape.shadow.inherit = False
    fill=shape.fill
    fill.solid()
    fill.fore_color.rgb=RGBColor(255,0,0)
    shape.text= "More people analytics"
    line=shape.line
    line.color.rgb=RGBColor(255,0,0)

    left = Inches(1)
    top = Inches(2)
    width = Inches(7)
    pic = slide.shapes.add_picture(imgpth, left, top, width=width)
    left = Inches(8)
    pic = slide.shapes.add_picture('hist.png', left, top, width=width)

    prs.save('typeform_presentation.pptx')

app = Flask(__name__)

@app.route('/downloadPPTX')
def plot_csv():
    generate_pptx_file()
    return send_file('typeform_presentation.pptx',
                     attachment_filename='typeform_presentation.pptx',
                     as_attachment=True)

if __name__ == '__main__':
    app.run()